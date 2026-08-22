import os
from pathlib import Path
import PyPDF2
import docx
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter


def extract_text_from_pdf(file_path: Path) -> str:
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        pages_text = [page.extract_text() for page in reader.pages]
        return "\n".join(t for t in pages_text if t)


def extract_text_from_docx(file_path: Path) -> str:
    doc = docx.Document(str(file_path))
    paras = [para.text for para in doc.paragraphs if para.text]
    table_cells = [
        cell.text
        for table in doc.tables
        for row in table.rows
        for cell in row.cells
        if cell.text
    ]
    return "\n".join(paras + table_cells)


def extract_text_from_pptx(file_path: Path) -> str:
    prs = Presentation(str(file_path))
    text_runs = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    ]
    return "\n".join(text_runs)


def extract_text_from_txt(file_path: Path) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    elif suffix in (".doc", ".docx"):
        return extract_text_from_docx(file_path)
    elif suffix in (".ppt", ".pptx"):
        return extract_text_from_pptx(file_path)
    elif suffix == ".txt":
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file format: {suffix}")


def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """Splits input text into optimized chunks and strips empty whitespace."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    raw_chunks = splitter.split_text(text)
    # Filter empty or tiny noise chunks for faster and cleaner embedding indexing
    return [c.strip() for c in raw_chunks if c and len(c.strip()) > 10]
